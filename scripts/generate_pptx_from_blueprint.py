#!/usr/bin/env python3
"""Generate a basic editable PPTX from slide_blueprint.json.

This script is a minimal reference implementation. It creates editable titles,
body text, tables, and simple process diagrams using PowerPoint native shapes.
It also writes speaker notes to a separate Markdown file because python-pptx
has no stable official API for speaker notes.

Usage:
    python scripts/generate_pptx_from_blueprint.py \
      --blueprint examples/sample_slide_blueprint.json \
      --output output/sample_presentation.pptx
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True


def add_bullets(slide, bullets: list[str], left=0.8, top=1.35, width=5.7, height=4.5) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    if not bullets:
        return
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)


def add_process_diagram(slide, diagram: dict[str, Any], left=6.9, top=1.55, width=5.4, height=3.0) -> None:
    nodes = diagram.get("nodes") or []
    if not nodes:
        return
    n = len(nodes)
    box_w = width / max(n, 1) * 0.82
    gap = width / max(n, 1) * 0.18
    y = top + height / 2 - 0.35
    prev = None
    for i, node in enumerate(nodes):
        x = left + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.7))
        shape.text = str(node)
        shape.text_frame.paragraphs[0].font.size = Pt(13)
        if prev is not None:
            # Use a right arrow shape between nodes. It remains editable.
            arrow_x = x - gap * 0.85
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x), Inches(y + 0.2), Inches(gap * 0.7), Inches(0.28))
            arrow.text = ""
        prev = shape


def add_table(slide, table: dict[str, Any], left=0.8, top=1.45, width=11.6, height=3.6) -> None:
    headers = table.get("headers") or []
    rows = table.get("rows") or []
    if not headers:
        return
    row_count = len(rows) + 1
    col_count = len(headers)
    tbl_shape = slide.shapes.add_table(row_count, col_count, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = tbl_shape.table
    for col, header in enumerate(headers):
        tbl.cell(0, col).text = str(header)
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row[:col_count]):
            tbl.cell(r, c).text = str(value)


def add_footer(slide, source_ids: list[str], slide_no: int) -> None:
    text = f"{slide_no}"
    if source_ids:
        text += "  |  " + ", ".join(source_ids)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)


def build_ppt(blueprint: list[dict[str, Any]], output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    notes_lines: list[str] = []

    for item in blueprint:
        slide = prs.slides.add_slide(blank_layout)
        slide_no = int(item.get("slide_no", len(prs.slides)))
        title = item.get("title", f"Slide {slide_no}")
        layout = item.get("layout_type", "content")
        content = item.get("visible_content", {}) or {}
        source_ids = item.get("source_ids", []) or []

        add_title(slide, title)

        if layout == "title":
            add_bullets(slide, content.get("bullets", []), left=1.2, top=2.2, width=10.8, height=2.5)
        elif content.get("table"):
            add_table(slide, content["table"])
        elif content.get("diagram") and content["diagram"].get("type") == "process":
            add_bullets(slide, content.get("bullets", []), left=0.8, top=1.45, width=5.5, height=4.6)
            add_process_diagram(slide, content["diagram"])
        else:
            add_bullets(slide, content.get("bullets", []), left=0.9, top=1.45, width=11.5, height=4.8)

        add_footer(slide, source_ids, slide_no)
        notes_lines.append(f"## Slide {slide_no}: {title}\n")
        notes_lines.append(f"{item.get('speaker_notes', '')}\n")
        if source_ids:
            notes_lines.append("来源：" + ", ".join(source_ids) + "\n")
        notes_lines.append("\n")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    notes_path = output.with_suffix(".speaker_notes.md")
    notes_path.write_text("".join(notes_lines), encoding="utf-8")
    return notes_path


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--blueprint", required=True, help="Path to slide_blueprint.json")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    args = parser.parse_args()

    blueprint_path = Path(args.blueprint)
    output_path = Path(args.output)
    blueprint = json.loads(blueprint_path.read_text(encoding="utf-8"))
    notes_path = build_ppt(blueprint, output_path)
    print(f"Wrote PPTX: {output_path}")
    print(f"Wrote notes: {notes_path}")


if __name__ == "__main__":
    main()
